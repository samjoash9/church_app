from fastapi import FastAPI, Depends, HTTPException, UploadFile, File, Form, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import StreamingResponse, FileResponse, JSONResponse
from sqlalchemy.orm import Session
from typing import List
import os
import io
import uuid

import models, schemas, database, ppt_themes

models.Base.metadata.create_all(bind=database.engine)

app = FastAPI(title="Worship Pads API")

# ---------------------------------------------------------------------------
# Bold font instancing (used by /api/export/ppt)
# ---------------------------------------------------------------------------
# Both bundled TTFs (TheSeasons.ttf, SpaceGrotesk-VariableFont_wght.ttf) are
# variable fonts with a "wght" axis rather than separate regular/bold files —
# same files the mobile app bundles, where Flutter's Skia renderer picks the
# Bold instance on the fly for `FontWeight.bold`. PIL/Pillow can't do that
# axis lookup itself, so we pre-bake a static Bold instance per font once
# (cached on disk) and hand PIL that file whenever bold text is requested.
_BOLD_FONT_CACHE: dict = {}


def _bold_variant_path(regular_path: str) -> str:
    """Returns the path to a static Bold (wght=700) instance of a variable
    font, instantiating and caching it on first use."""
    if regular_path in _BOLD_FONT_CACHE:
        return _BOLD_FONT_CACHE[regular_path]

    cache_dir = os.path.join(os.path.dirname(__file__), ".font_cache")
    os.makedirs(cache_dir, exist_ok=True)
    bold_path = os.path.join(cache_dir, os.path.basename(regular_path) + ".bold.ttf")

    if not os.path.exists(bold_path):
        from fontTools.ttLib import TTFont
        from fontTools.varLib.instancer import instantiateVariableFont

        font = TTFont(regular_path)
        if "fvar" in font:
            instantiateVariableFont(font, {"wght": 700}, inplace=True)
        font.save(bold_path)

    _BOLD_FONT_CACHE[regular_path] = bold_path
    return bold_path

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = os.environ.get("UPLOAD_DIR", os.path.join(os.environ.get("DB_DIR", "."), "uploads"))
os.makedirs(UPLOAD_DIR, exist_ok=True)

# ---------------------------------------------------------------------------
# Songs
# ---------------------------------------------------------------------------

@app.get("/api/songs", response_model=List[schemas.SongData])
def get_songs(db: Session = Depends(database.get_db)):
    return db.query(models.Song).all()

@app.post("/api/songs", response_model=schemas.SongData)
def create_song(song: schemas.SongDataCreate, db: Session = Depends(database.get_db)):
    db_song = db.query(models.Song).filter(models.Song.id == song.id).first()
    if db_song:
        raise HTTPException(status_code=400, detail="Song already exists")
    lines_data = [line.model_dump() for line in song.lines]
    new_song = models.Song(
        id=song.id,
        title=song.title,
        songKey=song.songKey,
        lines=lines_data,
        language=song.language,
    )
    db.add(new_song)
    db.commit()
    db.refresh(new_song)
    return new_song

@app.put("/api/songs/{song_id}", response_model=schemas.SongData)
def update_song(song_id: str, song: schemas.SongDataCreate, db: Session = Depends(database.get_db)):
    db_song = db.query(models.Song).filter(models.Song.id == song_id).first()
    if not db_song:
        raise HTTPException(status_code=404, detail="Song not found")
    db_song.title = song.title
    db_song.songKey = song.songKey
    db_song.lines = [line.model_dump() for line in song.lines]
    db_song.language = song.language
    db.commit()
    db.refresh(db_song)
    return db_song

@app.delete("/api/songs/{song_id}")
def delete_song(song_id: str, db: Session = Depends(database.get_db)):
    db_song = db.query(models.Song).filter(models.Song.id == song_id).first()
    if not db_song:
        raise HTTPException(status_code=404, detail="Song not found")
    db.delete(db_song)
    db.commit()
    return {"message": "deleted"}

# ---------------------------------------------------------------------------
# Lineup
# ---------------------------------------------------------------------------

@app.get("/api/lineup", response_model=List[schemas.LineupItem])
def get_lineup(db: Session = Depends(database.get_db)):
    return db.query(models.LineupEntry).order_by(models.LineupEntry.order_index).all()

@app.post("/api/lineup", response_model=schemas.LineupItem)
def add_to_lineup(item: schemas.LineupItemCreate, db: Session = Depends(database.get_db)):
    existing = db.query(models.LineupEntry).filter(models.LineupEntry.song_id == item.song_id).first()
    if existing:
        return existing
    max_order = db.query(models.LineupEntry).order_by(models.LineupEntry.order_index.desc()).first()
    next_order = (max_order.order_index + 1) if max_order else 0
    new_item = models.LineupEntry(order_index=next_order, song_id=item.song_id)
    db.add(new_item)
    db.commit()
    db.refresh(new_item)
    return new_item

@app.put("/api/lineup/reorder")
def reorder_lineup(body: schemas.LineupReorder, db: Session = Depends(database.get_db)):
    entries = db.query(models.LineupEntry).all()
    entry_map = {e.song_id: e for e in entries}
    for idx, song_id in enumerate(body.song_ids):
        if song_id in entry_map:
            entry_map[song_id].order_index = idx
    db.commit()
    return {"message": "reordered"}

@app.delete("/api/lineup/clear")
def clear_lineup(db: Session = Depends(database.get_db)):
    db.query(models.LineupEntry).delete()
    db.commit()
    return {"message": "cleared"}

@app.delete("/api/lineup/{item_id}")
def remove_from_lineup(item_id: int, db: Session = Depends(database.get_db)):
    db_item = db.query(models.LineupEntry).filter(models.LineupEntry.id == item_id).first()
    if not db_item:
        raise HTTPException(status_code=404, detail="Item not found")
    db.delete(db_item)
    db.commit()
    return {"message": "deleted"}

# ---------------------------------------------------------------------------
# PPT Presentations
# ---------------------------------------------------------------------------

@app.get("/api/ppts", response_model=List[schemas.PptData])
def get_ppts(db: Session = Depends(database.get_db)):
    return db.query(models.PptPresentation).all()

@app.post("/api/ppts", response_model=schemas.PptData)
def create_ppt(ppt: schemas.PptCreate, db: Session = Depends(database.get_db)):
    existing = db.query(models.PptPresentation).filter(models.PptPresentation.id == ppt.id).first()
    if existing:
        raise HTTPException(status_code=400, detail="PPT already exists")
    new_ppt = models.PptPresentation(id=ppt.id, title=ppt.title, song_ids=ppt.song_ids)
    db.add(new_ppt)
    db.commit()
    db.refresh(new_ppt)
    return new_ppt

@app.put("/api/ppts/{ppt_id}", response_model=schemas.PptData)
def update_ppt(ppt_id: str, ppt: schemas.PptCreate, db: Session = Depends(database.get_db)):
    db_ppt = db.query(models.PptPresentation).filter(models.PptPresentation.id == ppt_id).first()
    if not db_ppt:
        raise HTTPException(status_code=404, detail="PPT not found")
    db_ppt.title = ppt.title
    db_ppt.song_ids = ppt.song_ids
    db.commit()
    db.refresh(db_ppt)
    return db_ppt

@app.delete("/api/ppts/{ppt_id}")
def delete_ppt(ppt_id: str, db: Session = Depends(database.get_db)):
    db_ppt = db.query(models.PptPresentation).filter(models.PptPresentation.id == ppt_id).first()
    if not db_ppt:
        raise HTTPException(status_code=404, detail="PPT not found")
    db.delete(db_ppt)
    db.commit()
    return {"message": "deleted"}

# ---------------------------------------------------------------------------
# Worship Pads: sound library
# ---------------------------------------------------------------------------

@app.get("/api/sounds", response_model=List[schemas.SoundEntryData])
def get_sounds(db: Session = Depends(database.get_db)):
    return db.query(models.SoundEntry).all()

@app.post("/api/sounds/upload", response_model=schemas.SoundEntryData)
async def upload_sound(
    mode: str = Form(...),
    key: str = Form(...),
    file: UploadFile = File(...),
    db: Session = Depends(database.get_db),
):
    ext = os.path.splitext(file.filename or "")[1] or ".mp3"
    stored_name = f"{uuid.uuid4().hex}{ext}"
    dest_path = os.path.join(UPLOAD_DIR, stored_name)
    contents = await file.read()
    with open(dest_path, "wb") as f:
        f.write(contents)

    new_sound = models.SoundEntry(
        mode=mode,
        key=key,
        name=file.filename or stored_name,
        path=f"/media/uploads/{stored_name}",
        size_in_bytes=len(contents),
        is_asset=False,
        is_active=False,
    )
    db.add(new_sound)
    db.commit()
    db.refresh(new_sound)
    return new_sound

@app.delete("/api/sounds/{sound_id}")
def delete_sound(sound_id: int, db: Session = Depends(database.get_db)):
    sound = db.query(models.SoundEntry).filter(models.SoundEntry.id == sound_id).first()
    if not sound:
        raise HTTPException(status_code=404, detail="Sound not found")
    if not sound.is_asset:
        file_path = os.path.join(UPLOAD_DIR, os.path.basename(sound.path))
        if os.path.exists(file_path):
            os.remove(file_path)
    db.delete(sound)
    db.commit()
    return {"message": "deleted"}

@app.put("/api/sounds/{sound_id}/activate", response_model=schemas.SoundEntryData)
def activate_sound(sound_id: int, db: Session = Depends(database.get_db)):
    sound = db.query(models.SoundEntry).filter(models.SoundEntry.id == sound_id).first()
    if not sound:
        raise HTTPException(status_code=404, detail="Sound not found")
    db.query(models.SoundEntry).filter(
        models.SoundEntry.mode == sound.mode,
        models.SoundEntry.key == sound.key,
    ).update({"is_active": False})
    sound.is_active = True
    db.commit()
    db.refresh(sound)
    return sound

@app.put("/api/sounds/{mode}/{key}/clear-active")
def clear_active_sound(mode: str, key: str, db: Session = Depends(database.get_db)):
    db.query(models.SoundEntry).filter(
        models.SoundEntry.mode == mode,
        models.SoundEntry.key == key,
    ).update({"is_active": False})
    db.commit()
    return {"message": "cleared"}

# ---------------------------------------------------------------------------
# Export: PDF
# ---------------------------------------------------------------------------

@app.post("/api/export/pdf")
def export_pdf(req: schemas.ExportPdfRequest, db: Session = Depends(database.get_db)):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors as rl_colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Flowable
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.pdfbase.pdfmetrics import stringWidth

    class ChordLine(Flowable):
        """Chords drawn at their slot x-positions, mirroring the web editor's
        slot layout: uniform minimum slot width, slots grow to fit their chord,
        empty slots keep their width so later chords stay above the same spot.

        Web geometry (px): 20px min slot, 8px padding, 1px gap, 12px chord text
        above 14px lyric text. Everything scales by (PDF lyric pt / web lyric px)
        so each chord lands over the same word it covers on screen."""

        # 11/14 converts web px to PDF pt (11pt lyric vs 14px lyric); 0.87
        # compensates for Helvetica rendering narrower than the browser's
        # Segoe UI at equal nominal size, so slots shrink with the lyric.
        _SCALE = (11.0 / 14.0) * 0.87
        SLOT_MIN = 20 * _SCALE
        PAD = 8 * _SCALE
        GAP = 1 * _SCALE
        FONT_SIZE = 12 * _SCALE         # web chord text is 12px

        def __init__(self, chords, font="Helvetica-Bold", size=None, color=None):
            super().__init__()
            self.chords = chords
            self.font = font
            self.size = size if size is not None else self.FONT_SIZE
            self.color = color

        def wrap(self, availWidth, availHeight):
            self.width = availWidth
            self.height = self.size + 4
            return (self.width, self.height)

        def draw(self):
            self.canv.setFont(self.font, self.size)
            if self.color:
                self.canv.setFillColor(self.color)
            x = 0.0
            for chord in self.chords:
                text_w = stringWidth(chord, self.font, self.size) if chord else 0.0
                slot_w = max(self.SLOT_MIN, text_w + self.PAD)
                if chord:
                    self.canv.drawString(x + (slot_w - text_w) / 2, 2, chord)
                x += slot_w + self.GAP

    songs = []
    for sid in req.song_ids:
        s = db.query(models.Song).filter(models.Song.id == sid).first()
        if s:
            songs.append(s)

    if not songs:
        raise HTTPException(status_code=404, detail="No songs found for export")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle("SongTitle", parent=styles["Heading1"], fontSize=18, spaceAfter=4)
    key_style = ParagraphStyle("Key", parent=styles["Normal"], fontSize=11, textColor=rl_colors.grey, spaceAfter=12)
    section_style = ParagraphStyle("Section", parent=styles["Normal"], fontSize=11, textColor=rl_colors.HexColor("#6366f1"), fontName="Helvetica-Bold", spaceAfter=4)
    lyric_style = ParagraphStyle("Lyric", parent=styles["Normal"], fontSize=11, spaceAfter=6)

    story = []
    for song_idx, song in enumerate(songs):
        if song_idx > 0:
            story.append(Spacer(1, 0.5*cm))
        story.append(Paragraph(song.title, title_style))
        story.append(Paragraph(f"Key of {song.songKey}", key_style))

        lines = song.lines or []
        for line in lines:
            lyrics = line.get("lyrics", "") if isinstance(line, dict) else line.lyrics
            chords = line.get("chords", []) if isinstance(line, dict) else line.chords

            if not lyrics.strip() and all(not c for c in chords):
                story.append(Spacer(1, 0.3*cm))
                continue

            is_section = lyrics.strip().startswith("[") and lyrics.strip().endswith("]")
            if is_section:
                story.append(Spacer(1, 0.2*cm))
                story.append(Paragraph(lyrics.strip(), section_style))
                continue

            if any(c for c in chords):
                # Trim trailing empty slots; interior empties keep positioning.
                trimmed = list(chords)
                while trimmed and not trimmed[-1]:
                    trimmed.pop()
                story.append(ChordLine(trimmed, color=rl_colors.HexColor("#6366f1")))
            if lyrics.strip():
                story.append(Paragraph(lyrics, lyric_style))

        story.append(Spacer(1, 1*cm))

    doc.build(story)
    buf.seek(0)

    filename = "chord_charts.pdf" if len(songs) > 1 else f"{songs[0].title}.pdf"
    return StreamingResponse(
        buf,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )

# ---------------------------------------------------------------------------
# Export: PPTX
# ---------------------------------------------------------------------------

@app.get("/api/ppt-themes")
def list_ppt_themes():
    """Themes for the export picker; previews are served from /media/assets."""
    return [
        {
            "id": t.id,
            "displayName": t.display_name,
            "preview": f"/media/assets/{t.preview_asset}",
        }
        for t in ppt_themes.ALL_THEMES
    ]


@app.post("/api/export/ppt")
def export_ppt(req: schemas.ExportPptRequest, db: Session = Depends(database.get_db)):
    """Builds the same deck as the mobile exporter (ppt_export_service.dart):
    themed title slide, intro sections, per-song title + lyric slides, outro
    sections — over the shared assets/ppt_backgrounds images.

    Each slide is rasterized to a full-bleed JPEG (background + text baked in
    with the actual bundled TTF), exactly like the mobile app's
    captureFromWidget approach. A PPTX text run only *references* a font by
    name — if "The Seasons" / "Space Grotesk" isn't installed on whatever
    machine opens the file, PowerPoint/WPS silently substitutes a fallback
    (often a symbol font), producing garbled slides. Baking text into the
    image sidesteps font availability entirely.
    """
    import random
    from pptx import Presentation
    from pptx.util import Emu
    from PIL import Image, ImageDraw, ImageFont

    songs = []
    for sid in req.song_ids:
        s = db.query(models.Song).filter(models.Song.id == sid).first()
        if s:
            songs.append(s)

    if not songs:
        raise HTTPException(status_code=404, detail="No songs found for export")

    theme = ppt_themes.find_theme(req.theme or "cloud")

    # Mobile renders on a 1280x720 logical-px canvas; match that exactly so
    # theme font sizes/positions need no rescaling.
    CANVAS_W, CANVAS_H = 1280, 720

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.33 in (16:9)
    prs.slide_height = Emu(6858000)   # 7.5 in
    blank_layout = prs.slide_layouts[6]

    def asset_file(rel_path):
        path = os.path.join(assets_path, rel_path.replace("/", os.sep))
        return path if os.path.exists(path) else None

    def pick_background():
        if not theme.randomize_background or len(theme.background_assets) == 1:
            return theme.default_background
        return random.choice(theme.background_assets)

    font_path = asset_file(theme.font_file)

    def font(size_px, bold=True):
        # PIL needs an integer pixel size. Both bundled TTFs are variable
        # fonts, so bold text uses a pre-baked static Bold (wght=700)
        # instance rather than PIL's faux-bold, matching how Skia/Flutter
        # resolves FontWeight.bold on the same font files on mobile.
        path = font_path or asset_file("fonts/SpaceGrotesk-VariableFont_wght.ttf")
        if bold:
            path = _bold_variant_path(path)
        return ImageFont.truetype(path, int(size_px))

    def load_background(rel_path):
        path = asset_file(rel_path)
        if not path:
            return Image.new("RGB", (CANVAS_W, CANVAS_H), "black")
        img = Image.open(path).convert("RGB")
        # Cover-fit, same as the mobile FittedBox(fit: BoxFit.cover).
        src_ratio = img.width / img.height
        dst_ratio = CANVAS_W / CANVAS_H
        if src_ratio > dst_ratio:
            new_h = CANVAS_H
            new_w = int(new_h * src_ratio)
        else:
            new_w = CANVAS_W
            new_h = int(new_w / src_ratio)
        img = img.resize((new_w, new_h), Image.LANCZOS)
        left = (new_w - CANVAS_W) // 2
        top = (new_h - CANVAS_H) // 2
        return img.crop((left, top, left + CANVAS_W, top + CANVAS_H))

    def draw_centered_lines(draw, lines, center_y=None):
        """lines = [(text, size_px, color_hex, letter_spacing_px, gap_after_px, bold)].
        Stacks lines vertically, each horizontally centered, around center_y
        (defaults to canvas center)."""
        rendered = []
        total_h = 0
        for text, size_px, color_hex, spacing, gap_after, bold in lines:
            f = font(size_px, bold=bold)
            # Manual letter-spacing: measure/draw glyph-by-glyph.
            widths = [draw.textlength(ch, font=f) for ch in text]
            line_w = sum(widths) + spacing * max(0, len(text) - 1)
            ascent, descent = f.getmetrics()
            line_h = ascent + descent
            rendered.append((text, f, color_hex, spacing, widths, line_w, line_h, gap_after))
            total_h += line_h + gap_after
        total_h -= rendered[-1][7] if rendered else 0

        y = (center_y if center_y is not None else CANVAS_H / 2) - total_h / 2
        for text, f, color_hex, spacing, widths, line_w, line_h, gap_after in rendered:
            x = (CANVAS_W - line_w) / 2
            for ch, w in zip(text, widths):
                draw.text((x, y), ch, font=f, fill=f"#{color_hex}")
                x += w + spacing
            y += line_h + gap_after

    def wrap_text(draw, text, f, max_width):
        words = text.split(" ")
        lines, current = [], ""
        for word in words:
            candidate = f"{current} {word}".strip()
            if draw.textlength(candidate, font=f) <= max_width or not current:
                current = candidate
            else:
                lines.append(current)
                current = word
        if current:
            lines.append(current)
        return lines

    def render_slide(background_rel, block_lines=None, wrap_lines=None):
        """block_lines: fixed short lines (titles/section headers), centered.
        wrap_lines: long lyric lines that may need word-wrap, appended below
        block_lines, all centered together as one group."""
        img = load_background(background_rel)
        draw = ImageDraw.Draw(img)

        all_lines = list(block_lines or [])
        if wrap_lines:
            margin = 96
            for text, size_px, color_hex, gap_after, bold in wrap_lines:
                f = font(size_px, bold=bold)
                for wrapped in wrap_text(draw, text, f, CANVAS_W - margin * 2):
                    all_lines.append((wrapped, size_px, color_hex, 0, gap_after, bold))

        draw_centered_lines(draw, all_lines)

        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        buf.seek(0)
        return buf

    def add_slide_image(buf):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(buf, 0, 0, width=prs.slide_width, height=prs.slide_height)

    def add_section_slide(title):
        # Mobile: fontWeight: FontWeight.bold.
        lines = [(title, theme.section_font_size, theme.section_text_color, 8, 0, True)]
        add_slide_image(render_slide(pick_background(), block_lines=lines))

    # ── Main title slide ────────────────────────────────────────────────────
    if theme.show_title_overlay:
        # Mobile also draws a date top-right and a church icon; the bundled
        # themes all ship pre-rendered title art (overlay off), so this branch
        # only needs the text parts.
        lines = [
            # "B L E S S E D" has no fontWeight override on mobile (normal).
            (theme.main_title_part1, 56, theme.title_text_color, 24, 8, False),
            # "SUNDAY" is fontWeight: FontWeight.bold.
            (theme.main_title_part2, 200, theme.title_text_color, 0, 0, True),
        ]
        add_slide_image(render_slide(theme.resolved_title_background, block_lines=lines))
    else:
        add_slide_image(render_slide(theme.resolved_title_background))

    # ── Intro sections ──────────────────────────────────────────────────────
    for section in ppt_themes.INTRO_SECTIONS:
        add_section_slide(section)

    # ── Song slides ─────────────────────────────────────────────────────────
    for song in songs:
        # Mobile: fontWeight: FontWeight.bold.
        title_lines = [(song.title.upper(), theme.title_font_size, theme.title_text_color, 8, 0, True)]
        add_slide_image(render_slide(pick_background(), block_lines=title_lines))

        for group in ppt_themes.build_song_slides(song.title, song.lines or []):
            # Mobile: both the section label and lyric lines are
            # fontWeight: FontWeight.bold.
            block = [(group["title"], theme.lyrics_title_font_size, theme.title_text_color, 0, 24, True)]
            wrap = [(lyric, theme.lyrics_font_size, theme.lyrics_text_color, 16, True) for lyric in group["lyrics"]]
            add_slide_image(render_slide(pick_background(), block_lines=block, wrap_lines=wrap))

    # ── Outro sections ──────────────────────────────────────────────────────
    for section in ppt_themes.OUTRO_SECTIONS:
        add_section_slide(section)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = "presentation.pptx" if len(songs) > 1 else f"{songs[0].title}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )

# ---------------------------------------------------------------------------
# Static assets + frontend
# ---------------------------------------------------------------------------

assets_path = os.environ.get(
    "ASSETS_PATH",
    os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "assets")),
)
if os.path.exists(assets_path):
    app.mount("/media/assets", StaticFiles(directory=assets_path), name="media_assets")

if os.path.exists(UPLOAD_DIR):
    app.mount("/media/uploads", StaticFiles(directory=UPLOAD_DIR), name="media_uploads")

frontend_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "frontend", "dist"))
if os.path.exists(frontend_path):
    # Serve the built Vite bundle (JS/CSS under /assets/*, favicon, etc.).
    app.mount("/assets", StaticFiles(directory=os.path.join(frontend_path, "assets")), name="frontend_assets")

    index_file = os.path.join(frontend_path, "index.html")

    # SPA fallback: any non-API, non-media path returns index.html so React
    # Router can handle client-side routes (/chords, /pads, …) on hard refresh.
    # A bare StaticFiles(html=True) mount only serves "/" and 404s deep links.
    @app.get("/{full_path:path}")
    def spa_fallback(full_path: str, request: Request):
        if full_path.startswith(("api/", "media/")):
            return JSONResponse({"detail": "Not Found"}, status_code=404)
        candidate = os.path.join(frontend_path, full_path)
        if full_path and os.path.isfile(candidate):
            return FileResponse(candidate)
        return FileResponse(index_file)

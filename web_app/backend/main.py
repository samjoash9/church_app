from fastapi import FastAPI, Depends, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from typing import List
import os
import io

import models, schemas, database

models.Base.metadata.create_all(bind=database.engine)

app = FastAPI(title="Worship Pads API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---------------------------------------------------------------------------
# Songs
# ---------------------------------------------------------------------------

@app.get("/api/songs", response_model=List[schemas.SongData])
def get_songs(db: Session = Depends(database.get_db)):
    return db.query(models.Song).all()

@app.post("/api/songs", response_model=schemas.SongData)
def create_song(song: schemas.SongDataCreate, db: Session = Depends(database.get_db)):
    db_song = db.query(models.Song).filter(models.Song.id == song.id).first()
    if db_song:
        raise HTTPException(status_code=400, detail="Song already exists")
    lines_data = [line.model_dump() for line in song.lines]
    new_song = models.Song(
        id=song.id,
        title=song.title,
        songKey=song.songKey,
        lines=lines_data,
        language=song.language,
    )
    db.add(new_song)
    db.commit()
    db.refresh(new_song)
    return new_song

@app.put("/api/songs/{song_id}", response_model=schemas.SongData)
def update_song(song_id: str, song: schemas.SongDataCreate, db: Session = Depends(database.get_db)):
    db_song = db.query(models.Song).filter(models.Song.id == song_id).first()
    if not db_song:
        raise HTTPException(status_code=404, detail="Song not found")
    db_song.title = song.title
    db_song.songKey = song.songKey
    db_song.lines = [line.model_dump() for line in song.lines]
    db_song.language = song.language
    db.commit()
    db.refresh(db_song)
    return db_song

@app.delete("/api/songs/{song_id}")
def delete_song(song_id: str, db: Session = Depends(database.get_db)):
    db_song = db.query(models.Song).filter(models.Song.id == song_id).first()
    if not db_song:
        raise HTTPException(status_code=404, detail="Song not found")
    db.delete(db_song)
    db.commit()
    return {"message": "deleted"}

# ---------------------------------------------------------------------------
# Lineup
# ---------------------------------------------------------------------------

@app.get("/api/lineup", response_model=List[schemas.LineupItem])
def get_lineup(db: Session = Depends(database.get_db)):
    return db.query(models.LineupEntry).order_by(models.LineupEntry.order_index).all()

@app.post("/api/lineup", response_model=schemas.LineupItem)
def add_to_lineup(item: schemas.LineupItemCreate, db: Session = Depends(database.get_db)):
    existing = db.query(models.LineupEntry).filter(models.LineupEntry.song_id == item.song_id).first()
    if existing:
        return existing
    max_order = db.query(models.LineupEntry).order_by(models.LineupEntry.order_index.desc()).first()
    next_order = (max_order.order_index + 1) if max_order else 0
    new_item = models.LineupEntry(order_index=next_order, song_id=item.song_id)
    db.add(new_item)
    db.commit()
    db.refresh(new_item)
    return new_item

@app.put("/api/lineup/reorder")
def reorder_lineup(body: schemas.LineupReorder, db: Session = Depends(database.get_db)):
    entries = db.query(models.LineupEntry).all()
    entry_map = {e.song_id: e for e in entries}
    for idx, song_id in enumerate(body.song_ids):
        if song_id in entry_map:
            entry_map[song_id].order_index = idx
    db.commit()
    return {"message": "reordered"}

@app.delete("/api/lineup/clear")
def clear_lineup(db: Session = Depends(database.get_db)):
    db.query(models.LineupEntry).delete()
    db.commit()
    return {"message": "cleared"}

@app.delete("/api/lineup/{item_id}")
def remove_from_lineup(item_id: int, db: Session = Depends(database.get_db)):
    db_item = db.query(models.LineupEntry).filter(models.LineupEntry.id == item_id).first()
    if not db_item:
        raise HTTPException(status_code=404, detail="Item not found")
    db.delete(db_item)
    db.commit()
    return {"message": "deleted"}

# ---------------------------------------------------------------------------
# PPT Presentations
# ---------------------------------------------------------------------------

@app.get("/api/ppts", response_model=List[schemas.PptData])
def get_ppts(db: Session = Depends(database.get_db)):
    return db.query(models.PptPresentation).all()

@app.post("/api/ppts", response_model=schemas.PptData)
def create_ppt(ppt: schemas.PptCreate, db: Session = Depends(database.get_db)):
    existing = db.query(models.PptPresentation).filter(models.PptPresentation.id == ppt.id).first()
    if existing:
        raise HTTPException(status_code=400, detail="PPT already exists")
    new_ppt = models.PptPresentation(id=ppt.id, title=ppt.title, song_ids=ppt.song_ids)
    db.add(new_ppt)
    db.commit()
    db.refresh(new_ppt)
    return new_ppt

@app.put("/api/ppts/{ppt_id}", response_model=schemas.PptData)
def update_ppt(ppt_id: str, ppt: schemas.PptCreate, db: Session = Depends(database.get_db)):
    db_ppt = db.query(models.PptPresentation).filter(models.PptPresentation.id == ppt_id).first()
    if not db_ppt:
        raise HTTPException(status_code=404, detail="PPT not found")
    db_ppt.title = ppt.title
    db_ppt.song_ids = ppt.song_ids
    db.commit()
    db.refresh(db_ppt)
    return db_ppt

@app.delete("/api/ppts/{ppt_id}")
def delete_ppt(ppt_id: str, db: Session = Depends(database.get_db)):
    db_ppt = db.query(models.PptPresentation).filter(models.PptPresentation.id == ppt_id).first()
    if not db_ppt:
        raise HTTPException(status_code=404, detail="PPT not found")
    db.delete(db_ppt)
    db.commit()
    return {"message": "deleted"}

# ---------------------------------------------------------------------------
# Export: PDF
# ---------------------------------------------------------------------------

@app.post("/api/export/pdf")
def export_pdf(req: schemas.ExportPdfRequest, db: Session = Depends(database.get_db)):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors as rl_colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm

    songs = []
    for sid in req.song_ids:
        s = db.query(models.Song).filter(models.Song.id == sid).first()
        if s:
            songs.append(s)

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle("SongTitle", parent=styles["Heading1"], fontSize=18, spaceAfter=4)
    key_style = ParagraphStyle("Key", parent=styles["Normal"], fontSize=11, textColor=rl_colors.grey, spaceAfter=12)
    section_style = ParagraphStyle("Section", parent=styles["Normal"], fontSize=11, textColor=rl_colors.HexColor("#6366f1"), fontName="Helvetica-Bold", spaceAfter=4)
    chord_style = ParagraphStyle("Chord", parent=styles["Normal"], fontSize=10, textColor=rl_colors.HexColor("#6366f1"), fontName="Helvetica-Bold", spaceAfter=0)
    lyric_style = ParagraphStyle("Lyric", parent=styles["Normal"], fontSize=11, spaceAfter=6)

    story = []
    for song_idx, song in enumerate(songs):
        if song_idx > 0:
            story.append(Spacer(1, 0.5*cm))
        story.append(Paragraph(song.title, title_style))
        story.append(Paragraph(f"Key of {song.songKey}", key_style))

        lines = song.lines or []
        for line in lines:
            lyrics = line.get("lyrics", "") if isinstance(line, dict) else line.lyrics
            chords = line.get("chords", []) if isinstance(line, dict) else line.chords

            if not lyrics.strip() and all(not c for c in chords):
                story.append(Spacer(1, 0.3*cm))
                continue

            is_section = lyrics.strip().startswith("[") and lyrics.strip().endswith("]")
            if is_section:
                story.append(Spacer(1, 0.2*cm))
                story.append(Paragraph(lyrics.strip(), section_style))
                continue

            active_chords = [c for c in chords if c]
            if active_chords:
                story.append(Paragraph("  ".join(active_chords), chord_style))
            if lyrics.strip():
                story.append(Paragraph(lyrics, lyric_style))

        story.append(Spacer(1, 1*cm))

    doc.build(story)
    buf.seek(0)

    filename = "chord_charts.pdf" if len(songs) > 1 else f"{songs[0].title}.pdf"
    return StreamingResponse(
        buf,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )

# ---------------------------------------------------------------------------
# Export: PPTX
# ---------------------------------------------------------------------------

@app.post("/api/export/ppt")
def export_ppt(req: schemas.ExportPptRequest, db: Session = Depends(database.get_db)):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    songs = []
    for sid in req.song_ids:
        s = db.query(models.Song).filter(models.Song.id == sid).first()
        if s:
            songs.append(s)

    theme = req.theme or "dark"
    if theme == "dark":
        bg_color = RGBColor(0x0d, 0x0d, 0x1a)
        text_color = RGBColor(0xff, 0xff, 0xff)
        chord_color = RGBColor(0x99, 0x8f, 0xff)
    else:
        bg_color = RGBColor(0xff, 0xff, 0xff)
        text_color = RGBColor(0x11, 0x11, 0x22)
        chord_color = RGBColor(0x63, 0x66, 0xf1)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    def add_slide(song, lyric_lines):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Song title + key in top-left corner
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(5), Inches(0.5))
        tf = title_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{song.title}  •  Key of {song.songKey}"
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(0x88, 0x88, 0xaa)

        # Main lyrics block centered
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.8))
        tf = content_box.text_frame
        tf.word_wrap = True

        first = True
        for text in lyric_lines:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.size = Pt(36)
            run.font.color.rgb = text_color
            run.font.bold = False

    for song in songs:
        lines = song.lines or []
        # Group lyrics into slides of ~4 lyric lines each
        lyric_buffer = []
        for line in lines:
            lyrics = line.get("lyrics", "") if isinstance(line, dict) else line.lyrics
            if not lyrics.strip():
                continue
            is_section = lyrics.strip().startswith("[") and lyrics.strip().endswith("]")
            if is_section:
                if lyric_buffer:
                    add_slide(song, lyric_buffer)
                    lyric_buffer = []
                lyric_buffer = [lyrics.strip("[]")]
                continue
            lyric_buffer.append(lyrics)
            if len(lyric_buffer) >= 4:
                add_slide(song, lyric_buffer)
                lyric_buffer = []

        if lyric_buffer:
            add_slide(song, lyric_buffer)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = "presentation.pptx" if len(songs) > 1 else f"{songs[0].title}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )

# ---------------------------------------------------------------------------
# Static assets + frontend
# ---------------------------------------------------------------------------

assets_path = os.environ.get(
    "ASSETS_PATH",
    os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "assets")),
)
if os.path.exists(assets_path):
    # Mounted at /media (not /assets) to avoid shadowing the Vite build's
    # own /assets/* JS/CSS bundles served from the frontend dist mount below.
    app.mount("/media", StaticFiles(directory=assets_path), name="media")

frontend_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "frontend", "dist"))
if os.path.exists(frontend_path):
    app.mount("/", StaticFiles(directory=frontend_path, html=True), name="frontend")
